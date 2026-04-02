import logging
from typing import Optional, List

from fastapi import APIRouter, HTTPException, Header, UploadFile, File
from pydantic import BaseModel

from core.config import settings
from core.firebase_client import knowledge_manager as kb, firebase_storage as storage
from core.database import SessionLocal, Conversation, count_vector_knowledge

router = APIRouter(prefix="/knowledge", tags=["Knowledge"])
logger = logging.getLogger(__name__)


# ── Helpers Auth ──────────────────────────────────────────────────────────────

from core.auth import get_user_id_from_header


# ── Schemas ───────────────────────────────────────────────────────────────────

class KnowledgeCreateRequest(BaseModel):
    title: str
    content: str
    source: Optional[str] = ""
    type: Optional[str] = "text"  # "text" | "file" | "agent_added"


class KnowledgeUpdateRequest(BaseModel):
    additional_content: str
    append: Optional[bool] = True


class KnowledgeSearchRequest(BaseModel):
    query: str
    max_results: Optional[int] = 5


# ── Endpoints ─────────────────────────────────────────────────────────────────

@router.get("/")
async def list_knowledge(authorization: Optional[str] = Header(None)):
    """Liste tous les documents de connaissance de l'utilisateur sur Cloud SQL."""
    user_id = get_user_id_from_header(authorization)
    docs = kb.get_user_knowledge(user_id)
    return {"user_id": user_id, "count": len(docs), "documents": docs}


@router.post("/")
async def create_knowledge(
    body: KnowledgeCreateRequest,
    authorization: Optional[str] = Header(None),
):
    """Ajoute un document de connaissance vectorisé sur Cloud SQL."""
    user_id = get_user_id_from_header(authorization)
    if not body.title.strip() or not body.content.strip():
        raise HTTPException(status_code=400, detail="title et content sont requis")
        
    word_count = len(body.content.strip().split())
    
    doc_id = kb.add_knowledge(
        user_id=user_id,
        title=body.title.strip(),
        content=body.content.strip(),
        source=body.source or "manual",
        doc_type=body.type or "text",
        word_count=word_count
    )
    
    if not doc_id:
        raise HTTPException(status_code=500, detail="Erreur lors de la création du document vectoriel")

    logger.info(f"📚 Document vectorisé ajouté par {user_id} : {body.title} ({doc_id})")
    return {"id": doc_id, "title": body.title, "status": "created"}


@router.patch("/{doc_id}")
async def update_knowledge(
    doc_id: str,
    body: KnowledgeUpdateRequest,
    authorization: Optional[str] = Header(None),
):
    """Note: La mise à jour vectorielle nécessite de recalculer l'embedding."""
    user_id = get_user_id_from_header(authorization)
    # Pour l'instant on reste simple : on ne supporte pas l'update via PATCH atomique sur les vieux docs.
    # On encourage la création de nouveaux docs ou on implémentera la logique plus tard.
    raise HTTPException(status_code=501, detail="L'édition de documents vectoriels sera disponible prochainement. Veuillez créer un nouveau document.")


@router.delete("/{doc_id}")
async def delete_knowledge(doc_id: str, authorization: Optional[str] = Header(None)):
    """Supprime un document sur Cloud SQL."""
    user_id = get_user_id_from_header(authorization)
    # 1. Récupérer les infos du document (pour avoir le storage path si file_url existe)
    doc = kb.get_knowledge(user_id, doc_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Document introuvable")
    
    # 2. Si c'est un fichier stocké, on tente de le supprimer de Storage
    file_url = doc.get("file_url")
    if file_url:
        try:
            storage_path = f"users/{user_id}/knowledge/{doc.get('title')}"
            storage.delete_file(storage_path)
        except Exception as e:
            logger.warning(f"Impossible de supprimer le fichier storage pour {doc_id}: {e}")

    # 3. Supprimer de la base vectorielle
    ok = kb.delete_knowledge(user_id, doc_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Document introuvable")
    return {"id": doc_id, "status": "deleted"}


@router.post("/search")
async def search_knowledge_docs(
    body: KnowledgeSearchRequest,
    authorization: Optional[str] = Header(None),
):
    """Recherche sémantique via pgvector sur Cloud SQL."""
    user_id = get_user_id_from_header(authorization)
    results = kb.search_knowledge(user_id, body.query, body.max_results or 5)
    return {"query": body.query, "count": len(results), "results": results}


@router.post("/refresh")
async def refresh_knowledge(authorization: Optional[str] = Header(None)):
    """Retourne simplement le nombre actuel de documents KB pour rafraîchir l'UI."""
    user_id = get_user_id_from_header(authorization)
    return {"status": "ok", "count": count_vector_knowledge(user_id)}


async def _describe_image_for_kb(raw_data: bytes, mime_type: str, filename: str) -> str:
    """Utilise Gemini Vision pour décrire une image et la rendre recherchable dans la KB."""
    try:
        import base64
        import google.generativeai as genai
        genai.configure(api_key=settings.GEMINI_API_KEY)
        model = genai.GenerativeModel(settings.GEMINI_MODEL)
        b64 = base64.b64encode(raw_data).decode("utf-8")
        response = model.generate_content([
            {
                "parts": [
                    {"inline_data": {"mime_type": mime_type, "data": b64}},
                    {"text": "Décris en détail et avec précision le contenu de cette image. Inclus : les objets présents, les personnes si présentes, les textes visibles, les couleurs dominantes, le contexte général. Cette description sera utilisée pour retrouver l'image dans une base de connaissances."}
                ]
            }
        ])
        return response.text.strip()
    except Exception as e:
        logger.warning(f"Gemini Vision indisponible pour {filename}: {e}")
        return f"[Image : {filename}] — Description automatique indisponible."


@router.post("/upload")
async def upload_knowledge_file(
    file: UploadFile = File(...),
    authorization: Optional[str] = Header(None),
):
    """Upload et vectorise un fichier (PDF, DOCX, TXT, MD, CSV, JSON, images, PPTX, etc.)."""
    user_id = get_user_id_from_header(authorization)

    content = ""
    filename = file.filename or "unknown"
    fn_lower = filename.lower()
    doc_type = "file"

    try:
        raw_data = await file.read()

        if fn_lower.endswith((".txt", ".md")):
            content = raw_data.decode("utf-8", errors="replace")

        elif fn_lower.endswith(".pdf"):
            import io
            from pypdf import PdfReader
            pdf = PdfReader(io.BytesIO(raw_data))
            content = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())

        elif fn_lower.endswith(".docx"):
            import io
            import docx
            doc_obj = docx.Document(io.BytesIO(raw_data))
            content = "\n".join(para.text for para in doc_obj.paragraphs if para.text.strip())

        elif fn_lower.endswith((".csv", ".json")):
            content = raw_data.decode("utf-8", errors="replace")

        elif fn_lower.endswith((".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp")):
            mime = file.content_type or "image/jpeg"
            doc_type = "image"
            description = await _describe_image_for_kb(raw_data, mime, filename)
            content = f"[Image : {filename}]\n\n{description}"

        elif fn_lower.endswith(".pptx"):
            import io
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw_data))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                    if slide_content:
                        slides_text.append(f"Slide {i}:\n" + "\n".join(slide_content))
                content = "\n\n".join(slides_text)
            except ImportError:
                raise HTTPException(status_code=400, detail="python-pptx requis pour les fichiers PPTX. Contactez l'administrateur.")

        else:
            # Tenter UTF-8, sinon rejeter proprement
            try:
                content = raw_data.decode("utf-8")
            except UnicodeDecodeError:
                raise HTTPException(
                    status_code=400,
                    detail=f"Format non supporté pour '{filename}'. Formats acceptés : PDF, DOCX, TXT, MD, CSV, JSON, JPG, PNG, WEBP, PPTX."
                )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur d'extraction du fichier {filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la lecture du fichier : {e}")

    if not content.strip():
        raise HTTPException(status_code=400, detail="Le fichier ne contient pas de contenu extractible")

    # 1. Upload du fichier original sur Firebase Storage (non-bloquant)
    file_url = None
    try:
        storage_path = f"users/{user_id}/knowledge/{filename}"
        file_url = storage.upload_file(
            bucket_name="knowledge-base",
            path=storage_path,
            file_bytes=raw_data,
            content_type=file.content_type or "application/octet-stream"
        )
    except Exception as e:
        logger.warning(f"Upload Storage ignoré (non-bloquant): {e}")

    # 2. Vectorisation du contenu texte
    word_count = len(content.strip().split())

    doc_id = kb.add_knowledge(
        user_id=user_id,
        title=filename,
        content=content.strip(),
        source="upload",
        doc_type=doc_type,
        file_url=file_url,
        word_count=word_count
    )

    logger.info(f"📚 Fichier vectorisé par {user_id} : {filename} ({doc_id})")
    return {"id": doc_id, "title": filename, "status": "created", "content": content.strip()[:100] + "..."}
