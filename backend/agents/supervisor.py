"""
Supervisor Agent — FLARE AI.
Architecture Supervisor-Worker : route les requêtes vers le worker spécialisé,
puis synthétise la réponse finale pour l'utilisateur.

Le Supervisor ne possède AUCUN outil — il décide et orchestre.
Les Workers exécutent avec leurs propres outils et LLM.
"""
import asyncio
import logging
import json
import re
import operator
import uuid
from typing import TypedDict, Annotated, Sequence, Literal, Optional, List

from langchain_core.messages import BaseMessage, HumanMessage, AIMessage
from langchain_core.tools import tool
from langgraph.graph import StateGraph, END
from langchain_core.runnables import RunnableConfig

from core.llm_factory import get_llm
from core.memory import SessionMemory, CoreMemory
from core.config import settings
from core.context import (
    current_user_id, generated_images, current_request_id,
    current_session_id, knowledge_saved, GLOBAL_IMAGE_REGISTRY, core_memory,
    current_inline_file,
)

logger = logging.getLogger(__name__)

CHAT_SHORT_TERM_MEMORY_MESSAGES = 10


# ─── Classification des requêtes ─────────────────────────────────────────────

# Mots-clés pour routage rapide (avant même d'appeler le LLM)
_RESEARCH_KEYWORDS = {
    "cherche", "recherche", "trouve", "search", "google", "web", "actualité",
    "info", "information", "deep research", "analyse", "étude", "compare",
    "statistique", "chiffre", "source", "article", "news", "rappelle",
    "mémorise", "souviens", "retiens", "knowledge", "connaissance", "base de connaissances",
}
_MEDIA_KEYWORDS = {
    # Génération image
    "image", "photo", "illustration", "dessin", "génère une image", "crée une image",
    "fais une image", "fais moi une image", "fais-moi une image",
    "génère moi une image", "visualise", "photomontage", "infographie",
    "logo", "bannière", "affiche", "portrait", "peinture", "avatar", "cover",
    "visuel", "mockup", "flyer", "poster", "icône", "thumbnail", "miniature",
    # Recréation / transformation image
    "recréer", "recrée", "recreer", "recree", "reproduis", "reproduire",
    "transforme cette image", "transforme la photo", "mets dans", "place dans",
    "photomonte", "photomanipulation",
    # Fond
    "fond transparent", "supprimer le fond", "enlever le fond", "retirer le fond",
    "découper", "détourage", "détourer", "changer le fond", "changer le décor",
    "changer l'arrière-plan", "nouvel arrière-plan", "nouveau fond", "nouveau décor",
    "mettre dans un bureau", "placer dans", "arrière-plan",
    # Vidéo
    "vidéo", "video", "clip", "animation", "anime", "animer", "animé", "veo", "imagen",
    "met en mouvement", "faire bouger", "mettre en mouvement", "donne vie",
    # Autres
    "skill", "compétence", "template",
}
_WORKSPACE_KEYWORDS = {
    "drive", "dossier",
    "gmail", "email", "mail", "envoie", "calendrier",
    "calendar", "agenda", "événement", "rdv", "rendez-vous", "meeting",
    "prospection", "campagne", "facebook", "messenger",
}
_DOCUMENT_KEYWORDS = {
    # Formats
    "word", "docx", "document word", "fichier word",
    # Actions
    "rédige", "rédiger", "génère un document", "crée un document", "fais un document",
    "génère un rapport", "crée un rapport", "fais un rapport",
    # Types de documents
    "rapport", "lettre", "cv", "rédaction", "guide", "manuel",
    "cours", "formation", "tutoriel", "leçon", "programme", "syllabus",
    "contrat", "devis", "facture word", "proposition",
    "fiche", "notice", "procédure", "cahier des charges",
    "compte-rendu", "compte rendu", "procès-verbal", "bilan",
    "essai", "mémoire", "thèse", "dissertation",
    "plan de", "note de synthèse", "synthèse",
}
_SHEET_KEYWORDS = {
    # Formats
    "excel", "xlsx", "tableur", "spreadsheet", "feuille de calcul", "csv",
    # Actions
    "génère un tableau", "crée un tableau", "fais un tableau",
    "génère un excel", "crée un excel", "fais un excel",
    # Types
    "budget", "facture", "planning", "plan de travail",
    "suivi", "tracker", "inventaire", "stock",
    "données", "base de données", "tableau de bord", "dashboard",
    "calcul", "formule", "statistiques",
}

def classify_request(text: str, model_override: Optional[str] = None,
                     file_type: Optional[str] = None, file_name: Optional[str] = None) -> str:
    """Classification ultra-rapide par mots-clés uniquement (<1ms, 0 token LLM).

    L'appel LLM de routage a été supprimé car il ajoutait 2-5s de latence
    pour un gain de précision négligeable. Les mots-clés couvrent 95%+ des cas.

    Si un fichier est joint, la classification tient compte du type de fichier
    pour router vers le bon worker ou vers "chat" pour l'analyse.
    """
    lower = text.lower()
    fname_lower = (file_name or "").lower()

    # Fichier joint sans mot-clé de génération → analyse en mode chat
    # (sauf si le texte indique explicitement une action de création)
    _has_generation_intent = any(kw in lower for kw in [
        "génère", "crée", "fais", "rédige", "produis", "écris",
        "generate", "create", "make", "write",
    ])
    if file_type and not _has_generation_intent:
        # Image ou vidéo jointe → chat multimodal (analyse)
        if file_type.startswith("image/") or file_type.startswith("video/"):
            # Sauf si l'utilisateur veut explicitement transformer/animer l'image
            _media_action = any(kw in lower for kw in [
                "anime", "animer", "transforme", "fond transparent", "supprime le fond",
                "enlève le fond", "retire le fond", "change le fond", "détoure",
                "met en mouvement", "donne vie", "fais bouger",
            ])
            if _media_action:
                return "media"
            logger.info(f"[Supervisor] File analysis mode: chat (file={file_type})")
            return "chat"
        # Document joint → chat pour analyse (pas génération)
        if any(fname_lower.endswith(ext) for ext in (".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".csv", ".txt")):
            logger.info(f"[Supervisor] Document analysis mode: chat (file={file_name})")
            return "chat"

    # 1. Salutations simples → chat immédiat
    greetings = r"^(bonjour|salut|hello|hi|hey|coucou|bonsoir|hola|yo)(\s+(flair|ia|ai|l'ia|l'ai|kévin|kevin))?[\s!?.1]*$"
    if re.match(greetings, lower):
        logger.info("[Supervisor] Fast Path: chat (greeting)")
        return "chat"

    # 2. Mots-clés critiques (match exact rapide)
    if any(kw in lower for kw in [
        # Image
        "génère une image", "génère moi une image", "crée une image", "fais une image",
        "dessine", "génère un logo", "crée un logo", "fais un logo", "fais moi un logo",
        "génère une affiche", "fais une affiche", "génère un visuel", "crée un visuel",
        "génère un avatar", "génère une illustration", "génère un portrait",
        "recréer", "recrée", "recree",
        # Fond
        "supprime le fond", "enlève le fond", "retire le fond", "fond transparent",
        "découpe cette image", "découpe la photo", "détoure", "détourage",
        "change le fond", "change le décor", "change l'arrière-plan",
        "mets dans un bureau", "mets cette personne dans", "place dans un",
        # Vidéo
        "génère une vidéo", "genere une vidéo", "génère une video", "genere une video",
        "crée une vidéo", "cree une vidéo", "crée une video", "cree une video",
        "fais une vidéo", "fais une video", "fais moi une vidéo", "fais moi une video",
        "fais-moi une vidéo", "fais-moi une video", "créer une vidéo", "creer une video",
        "génère moi une vidéo", "genere moi une video",
        "anime cette image", "anime la photo", "anime cette photo", "anime l'image",
        "animer cette image", "animer la photo", "animer cette photo",
        "met en mouvement", "mettre en mouvement", "donne vie à",
        "fais bouger", "faire bouger",
    ]):
        return "media"

    if any(kw in lower for kw in [
        # Document Word
        "génère un document", "crée un document", "fais un document", "fais moi un document",
        "génère un rapport", "crée un rapport", "fais un rapport", "fais moi un rapport",
        "génère un cours", "crée un cours", "fais un cours", "fais moi un cours",
        "génère un cv", "crée un cv", "fais un cv",
        "génère une lettre", "crée une lettre", "fais une lettre",
        "génère un guide", "crée un guide", "fais un guide",
        "génère un contrat", "crée un contrat", "fais un contrat",
        "rédige un", "rédiger un",
    ]):
        return "document"

    if any(kw in lower for kw in [
        # Excel
        "génère un excel", "crée un excel", "fais un excel", "fais moi un excel",
        "génère un tableur", "crée un tableur", "fais un tableur",
        "génère un budget", "crée un budget", "fais un budget",
        "génère un planning", "crée un planning", "fais un planning",
        "génère un tableau", "crée un tableau", "fais un tableau de",
        "génère une facture", "crée une facture", "fais une facture",
    ]):
        return "sheet"

    if lower.startswith("/image") or lower.startswith("/video"):
        return "media"
    if lower.startswith("/doc") or lower.startswith("/word"):
        return "document"
    if lower.startswith("/sheet") or lower.startswith("/excel"):
        return "sheet"
    if any(kw in lower for kw in ["cherche sur le web", "deep research", "recherche approfondie"]):
        return "researcher"

    # 3. Scoring par mots-clés (pondéré — mots longs = plus pertinents)
    scores = {"researcher": 0, "media": 0, "workspace": 0, "document": 0, "sheet": 0}
    for kw in _RESEARCH_KEYWORDS:
        if kw in lower:
            scores["researcher"] += 2 if len(kw) > 6 else 1
    for kw in _MEDIA_KEYWORDS:
        if kw in lower:
            scores["media"] += 2 if len(kw) > 6 else 1
    for kw in _WORKSPACE_KEYWORDS:
        if kw in lower:
            scores["workspace"] += 2 if len(kw) > 6 else 1
    for kw in _DOCUMENT_KEYWORDS:
        if kw in lower:
            scores["document"] += 2 if len(kw) > 6 else 1
    for kw in _SHEET_KEYWORDS:
        if kw in lower:
            scores["sheet"] += 2 if len(kw) > 6 else 1

    max_score = max(scores.values())
    if max_score == 0:
        return "chat"

    for worker in ["researcher", "media", "workspace", "document", "sheet"]:
        if scores[worker] == max_score:
            logger.info(f"[Supervisor] Keyword Routing: {worker} (score={max_score})")
            return worker
    return "chat"


# ─── Nettoyage des réponses ──────────────────────────────────────────────────

def _clean_response(content) -> str:
    """Nettoie la réponse en extrayant le texte utile."""
    if not content:
        return ""
    if isinstance(content, list):
        text_parts = []
        for part in content:
            if isinstance(part, str):
                text_parts.append(part)
            elif hasattr(part, "text") and part.text:
                text_parts.append(part.text)
            elif isinstance(part, dict) and part.get("text"):
                text_parts.append(part["text"])
        content = "\n".join(text_parts)
    if not isinstance(content, str):
        content = str(content)
    if not content:
        return ""
    # Supprimer les blocs <thought>
    without_thoughts = re.sub(r'<thought>[\s\S]*?</thought>', '', content).strip()
    if without_thoughts:
        content = without_thoughts
    else:
        thought_texts = re.findall(r'<thought>([\s\S]*?)</thought>', content)
        if thought_texts:
            content = "\n".join(thought_texts).strip()
    content = re.sub(
        r'^\s*\{["\s]*name["\s]*:\s*"[a-z_]+"\s*,[^}]*"parameters"\s*:\s*\{[^}]*\}\s*\}\s*$',
        '', content, flags=re.MULTILINE
    )
    content = re.sub(r'\n{3,}', '\n\n', content)
    return content.strip()

def _extract_suggestions(text: str) -> tuple[str, list[str]]:
    """Extrait les suggestions [SUGGESTION: ...] du texte et retourne (texte_nettoyé, liste_suggestions)."""
    if not text:
        return "", []
    suggestions = []
    pattern = r'\[SUGGESTION:\s*(.*?)\]'
    matches = re.findall(pattern, text, re.IGNORECASE)
    for m in matches:
        s = m.strip()
        if s and s not in suggestions:
            suggestions.append(s)
    cleaned_text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
    return cleaned_text, suggestions


def _looks_like_invalid_media_markdown(text: str) -> bool:
    if not text:
        return False
    stripped = text.strip()
    if "flare-model-hosting" in stripped:
        return True
    if re.search(r'!\[[^\]]*\]\((https?://[^\)]+)\)', stripped) and "storage.googleapis.com" in stripped:
        return True
    if re.search(r'^\s*\|.*\|\s*$', stripped, re.MULTILINE) and "![" in stripped and "storage.googleapis.com" in stripped:
        return True
    return False


# ─── Prompt système du Supervisor ─────────────────────────────────────────────

SUPERVISOR_SYSTEM_PROMPT = """Tu es FLARE AI. Date du jour : {current_date}.

- Tu es l’assistant personnel de l’utilisateur dans FLARE AI.
- Tu aides l’utilisateur à comprendre, décider, créer, analyser et agir.
- Réponds avec clarté, simplicité, précision et honnêteté.
- N’invente jamais de faits, de sources, de chiffres, d’actualités ou de capacités.
- N’affirme jamais avoir utilisé un outil si ce n’est pas vrai.
- Utilise la recherche web pour les informations récentes ou à vérifier.
- Utilise la mémoire et la base de connaissances seulement si elles aident vraiment.
- Tu peux converser, rechercher, analyser des fichiers, créer ou modifier Word et Excel, générer ou retoucher des images, générer des vidéos et animer une image si une image source est disponible.
- Ne parle pas de worker, de module, de prompt ou de détails internes, sauf si l’utilisateur le demande.
- Si on te demande ce que tu es, réponds simplement : "Je suis FLARE AI."
"""


# ─── Noms d'outils lisibles (pour le frontend) ──────────────────────────────

SUPERVISOR_RUNTIME_PROMPT = """Tu es FLARE AI. Date du jour : {current_date}.

- Tu es l’assistant personnel de l’utilisateur dans FLARE AI.
- Français par défaut, sauf si l’utilisateur choisit une autre langue.
- Réponds de façon claire, simple, humaine, professionnelle et utile.
- Pas de blabla, pas de jargon inutile, pas de ton marketing, pas de détails internes.
- N’invente rien : ni source, ni chiffre, ni actualité, ni capacité.
- Si une information dépend du temps et n’a pas été vérifiée, dis-le clairement.
- Si une action n’est pas disponible, dis-le simplement et propose la meilleure alternative.
- Quand la demande est claire et faisable, agis sans faire perdre de temps.
- Si l’utilisateur demande un média, un document ou un tableur, privilégie l’exécution.
- Si l’utilisateur fournit une image et demande de l’animer ou de la modifier, utilise cette image en priorité.
- Utilise la recherche web pour les faits récents ou à vérifier.
- Utilise la mémoire pour les faits durables et utiles, pas pour les détails jetables.
- Utilise la base de connaissances seulement si elle aide vraiment.
- Tu peux converser, rechercher, analyser des fichiers, créer ou modifier Word et Excel, générer ou retoucher des images, générer des vidéos et animer une image si une image source est disponible.
- Pour une simple salutation, salue puis demande l’objectif.
- Si on te demande ce que tu es, réponds simplement : "Je suis FLARE AI."
- Réponse complète mais compacte. Donne des étapes ou un exemple court seulement si utile.
- Quand tu cites des informations externes vérifiées, ajoute des sources sous la forme [titre](url).
- Ajoute exactement 2 suggestions courtes et actionnables à la fin quand elles sont utiles. Format strict :
[SUGGESTION: action concrète 1]
[SUGGESTION: action concrète 2]

{user_persona}
{core_memory}"""


def _is_pure_greeting(text: str) -> bool:
    if not text:
        return False
    return bool(re.match(r"^\s*(bonjour|salut|hello|hi|hey|bonsoir|coucou|yo)(\s+(flare|flare ai|ia|ai))?[\s!?.]*$", text.strip(), re.IGNORECASE))


def _build_greeting_response() -> tuple[str, list[str]]:
    return (
        "Bonjour. Je peux rechercher des informations, analyser un fichier, creer ou modifier un document Word, creer un tableau Excel, generer ou retoucher une image, generer une video, ou utiliser votre memoire et vos connaissances. Que voulez-vous faire ?",
        ["Lancer une recherche web", "Creer une image"],
    )


def _ensure_suggestions(suggestions: list[str], worker_type: str) -> list[str]:
    cleaned = []
    for suggestion in suggestions:
        text = (suggestion or "").strip()
        if text and text not in cleaned:
            cleaned.append(text)
    defaults = {
        "researcher": ["Creer un resume executif", "Comparer plusieurs options"],
        "media": ["Creer une variante plus premium", "Adapter au format mobile"],
        "document": ["Ameliorer la structure du document", "Generer une version plus concise"],
        "sheet": ["Ajouter des formules utiles", "Creer un tableau de bord visuel"],
        "workspace": ["Preparer le message a envoyer", "Structurer les prochaines actions"],
        "chat": ["Approfondir ce sujet", "Passer a une action concrete"],
    }
    for fallback in defaults.get(worker_type, defaults["chat"]):
        if len(cleaned) >= 2:
            break
        if fallback not in cleaned:
            cleaned.append(fallback)
    return cleaned[:2]


TOOL_DISPLAY_NAMES = {
    "web_search": "Recherche sur le web",
    "execute_deep_research": "Recherche approfondie en cours",
    "search_knowledge_base": "Consultation de la base de connaissances",
    "recall_facts": "Récupération des souvenirs utiles",
    "remember_fact": "Enregistrement d’une nouvelle information",
    "generate_image": "Création de l’image",
    "generate_video": "Création de la vidéo",
    "animate_image": "Animation de l’image",
    "edit_image_zone": "Retouche ciblée de l’image",
    "edit_video_clip": "Montage et modification de la vidéo",
    "generate_word_document": "Création du document Word",
    "read_word_document": "Lecture du document Word",
    "read_word_document_as_json": "Analyse de la structure du document Word",
    "update_word_document": "Mise à jour du document Word",
    "delete_word_document": "Suppression du document Word",
    "generate_excel_document": "Création du tableau Excel",
    "read_excel_document": "Lecture du tableau Excel",
    "read_excel_document_as_json": "Analyse de la structure du tableau Excel",
    "update_excel_document": "Mise à jour du tableau Excel",
    "delete_excel_document": "Suppression du tableau Excel",
}


# ─── Supervisor Agent ─────────────────────────────────────────────────────────

class SupervisorAgent:
    """
    Supervisor : route les requêtes vers les workers spécialisés.

    Architecture :
    - Classification rapide par mots-clés (0 token LLM)
    - Si "chat" (conversation simple) → le Supervisor répond directement
    - Si worker identifié → délègue au worker, puis synthétise si nécessaire
    """

    def __init__(self):
        # Le Supervisor utilise le modèle "cerveau" pour la conversation directe
        self.llm = get_llm(temperature=0.7)

        # Workers (lazy init pour ne pas charger tous les outils au démarrage)
        self._researcher = None
        self._media = None
        self._workspace = None
        self._document = None
        self._sheet = None

        logger.info("[SupervisorAgent] Initialisé — architecture Supervisor-Worker active")

    @property
    def researcher(self):
        if self._researcher is None:
            from agents.workers.researcher import ResearcherWorker
            self._researcher = ResearcherWorker()
        return self._researcher

    @property
    def media(self):
        if self._media is None:
            from agents.workers.media import MediaWorker
            self._media = MediaWorker()
        return self._media

    @property
    def workspace(self):
        if self._workspace is None:
            from agents.workers.workspace import WorkspaceWorker
            self._workspace = WorkspaceWorker()
        return self._workspace

    @property
    def document(self):
        if self._document is None:
            from agents.workers.document_worker import DocumentWorker
            self._document = DocumentWorker()
        return self._document

    @property
    def sheet(self):
        if self._sheet is None:
            from agents.workers.spreadsheet_worker import SpreadsheetWorker
            self._sheet = SpreadsheetWorker()
        return self._sheet

    def _build_human_message(
        self, text: str, file_content: Optional[str] = None,
        file_type: Optional[str] = None, file_name: Optional[str] = None,
    ) -> HumanMessage:
        """Construit le message utilisateur avec fichier joint et contexte de sélection."""
        
        prompt = text
        selection_context = ""

        # TKT-041: Essayer de parser le texte comme du JSON pour extraire la sélection
        try:
            data = json.loads(text)
            if isinstance(data, dict) and 'prompt' in data and 'selection' in data:
                prompt = data.get('prompt', '')
                selection = data.get('selection', {})
                sel_type = selection.get('type')
                
                if sel_type == 'text_selection' and 'selected_text' in selection:
                    selection_context = f"[CONTEXTE DE SÉLECTION]\nFichier: {selection.get('file_url', 'inconnu')}\nTexte sélectionné: \"{selection['selected_text']}\"\n\n"
                elif sel_type == 'cell_selection' and 'range' in selection:
                    selection_context = f"[CONTEXTE DE SÉLECTION]\nFichier: {selection.get('file_url', 'inconnu')}\nPlage sélectionnée: {selection.get('sheet_name', '')}!{selection['range']}\n\n"
                elif sel_type == 'document_refinement':
                    file_url = selection.get('file_url', '')
                    file_name = selection.get('file_name', '') or (file_name or '')
                    if isinstance(file_url, str) and file_url.startswith("data:"):
                        file_url = ""
                    if not file_url and file_content and file_name:
                        file_url = "inline://current-file"
                    selection_context = f"[MODIFICATION FICHIER]\nURL du fichier a modifier: {file_url}\nNom du fichier: {file_name}\n\n"
                elif sel_type == 'image_refinement':
                    image_url = selection.get('image_url', '')
                    image_name = selection.get('image_name', '') or (file_name or '')
                    action = selection.get('action', '')
                    extra = selection.get('extra', '')
                    mask_attached = "oui" if file_content and (file_type or "").startswith("image/") else "non"
                    selection_context = (
                        f"[MODIFICATION IMAGE]\n"
                        f"Action demandee: {action}\n"
                        f"URL de l'image a modifier: {image_url}\n"
                        f"Nom de l'image: {image_name}\n"
                        f"Masque inline attache: {mask_attached}\n"
                        f"Contexte supplementaire: {extra}\n\n"
                    )

                # Le prompt final pour le worker inclura ce contexte
                text = f"{selection_context}{prompt}"

        except (json.JSONDecodeError, TypeError):
            # Ce n'est pas du JSON, on continue normalement
            pass

        if not file_content or not file_type:
            return HumanMessage(content=text)

        fname_lower = (file_name or "").lower()

        # ── Images — multimodal natif (Gemini, VertexAI, OpenAI) ──────────
        if file_type.startswith("image/"):
            if settings.LLM_PROVIDER in ("openai", "gemini", "vertexai"):
                return HumanMessage(content=[
                    {"type": "text", "text": text or "Analyse cette image en détail. Décris ce que tu vois, le contexte, les éléments importants."},
                    {"type": "image_url", "image_url": {"url": f"data:{file_type};base64,{file_content}"}}
                ])
            return HumanMessage(content=f"[Image jointe : {file_name}]\n{text}")

        # ── Vidéos — multimodal natif Gemini (analyse de frames) ──────────
        elif file_type.startswith("video/"):
            if settings.LLM_PROVIDER in ("gemini", "vertexai"):
                return HumanMessage(content=[
                    {"type": "text", "text": text or "Analyse cette vidéo en détail. Décris ce que tu vois, les actions, le contexte, les éléments importants."},
                    {"type": "image_url", "image_url": {"url": f"data:{file_type};base64,{file_content}"}}
                ])
            return HumanMessage(content=f"[Vidéo jointe : {file_name}]\nL'analyse vidéo nécessite le mode Gemini.\n{text}")

        # ── Audio ─────────────────────────────────────────────────────────
        elif file_type.startswith("audio/"):
            if not text or "[Message vocal]" not in text:
                return HumanMessage(content=f"[Message vocal : {file_name}]\nTranscription échouée.\n{text or ''}")
            return HumanMessage(content=text)

        # ── PDF ───────────────────────────────────────────────────────────
        elif fname_lower.endswith(".pdf"):
            try:
                import base64 as b64mod, io
                from pypdf import PdfReader
                raw = b64mod.b64decode(file_content)
                pdf = PdfReader(io.BytesIO(raw))
                pages_text = []
                for i, p in enumerate(pdf.pages):
                    page_text = p.extract_text()
                    if page_text:
                        pages_text.append(f"--- Page {i+1} ---\n{page_text}")
                extracted = "\n\n".join(pages_text)
                preview = extracted[:12000] + ("..." if len(extracted) > 12000 else "")
                meta_info = f"({len(pdf.pages)} pages)"
                return HumanMessage(content=f"[PDF : {file_name} {meta_info}]\n\n{preview}\n\n{text}")
            except Exception as e:
                return HumanMessage(content=f"[PDF : {file_name} — extraction échouée : {e}]\n\n{text}")

        # ── DOCX (Word) ──────────────────────────────────────────────────
        elif fname_lower.endswith(".docx"):
            try:
                import base64 as b64mod, io, docx as _docx
                raw = b64mod.b64decode(file_content)
                doc_obj = _docx.Document(io.BytesIO(raw))
                parts = []
                for para in doc_obj.paragraphs:
                    if para.text.strip():
                        # Préserver les titres
                        if para.style and para.style.name and para.style.name.startswith("Heading"):
                            level = para.style.name.replace("Heading", "").strip() or "1"
                            parts.append(f"{'#' * int(level)} {para.text}")
                        else:
                            parts.append(para.text)
                # Extraire les tableaux
                for table in doc_obj.tables:
                    parts.append("\n[Tableau]")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))
                extracted = "\n".join(parts)
                preview = extracted[:12000] + ("..." if len(extracted) > 12000 else "")
                return HumanMessage(content=f"[DOCX : {file_name}]\n\n{preview}\n\n{text}")
            except Exception as e:
                return HumanMessage(content=f"[DOCX : {file_name} — erreur : {e}]\n\n{text}")

        # ── XLSX (Excel) ─────────────────────────────────────────────────
        elif fname_lower.endswith(".xlsx") or fname_lower.endswith(".xls"):
            try:
                import base64 as b64mod, io
                from openpyxl import load_workbook
                raw = b64mod.b64decode(file_content)
                wb = load_workbook(io.BytesIO(raw), data_only=True)
                extracted = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    extracted.append(f"--- Feuille: {sheet_name} ({ws.max_row} lignes x {ws.max_column} colonnes) ---")
                    row_count = 0
                    for row in ws.iter_rows(values_only=True):
                        if not any(row):
                            continue
                        extracted.append(" | ".join([str(c) if c is not None else "" for c in row]))
                        row_count += 1
                        if row_count >= 200:  # Limiter pour les gros fichiers
                            extracted.append(f"... ({ws.max_row - 200} lignes supplémentaires)")
                            break

                preview = "\n".join(extracted)
                preview = preview[:12000] + ("..." if len(preview) > 12000 else "")
                return HumanMessage(content=f"[TABLEUR : {file_name} — {len(wb.sheetnames)} feuille(s)]\n\n{preview}\n\n{text}")
            except Exception as e:
                return HumanMessage(content=f"[TABLEUR : {file_name} — erreur : {e}]\n\n{text}")

        # ── PPTX (PowerPoint) ────────────────────────────────────────────
        elif fname_lower.endswith(".pptx"):
            try:
                import base64 as b64mod, io
                from pptx import Presentation
                raw = b64mod.b64decode(file_content)
                prs = Presentation(io.BytesIO(raw))
                parts = []
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    slide_texts.append(t)
                        if shape.has_table:
                            for row in shape.table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                slide_texts.append(" | ".join(cells))
                    if slide_texts:
                        parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))
                extracted = "\n\n".join(parts)
                preview = extracted[:12000] + ("..." if len(extracted) > 12000 else "")
                return HumanMessage(content=f"[PPTX : {file_name} — {len(prs.slides)} slides]\n\n{preview}\n\n{text}")
            except Exception as e:
                return HumanMessage(content=f"[PPTX : {file_name} — erreur : {e}]\n\n{text}")

        # ── CSV ──────────────────────────────────────────────────────────
        elif fname_lower.endswith(".csv"):
            try:
                import base64 as b64mod
                raw = b64mod.b64decode(file_content)
                # Essayer plusieurs encodages
                for enc in ("utf-8", "latin-1", "cp1252"):
                    try:
                        decoded = raw.decode(enc)
                        break
                    except UnicodeDecodeError:
                        decoded = raw.decode("utf-8", errors="replace")
                line_count = decoded.count("\n")
                preview = decoded[:12000] + ("..." if len(decoded) > 12000 else "")
                return HumanMessage(content=f"[CSV : {file_name} — ~{line_count} lignes]\n\n{preview}\n\n{text}")
            except Exception as e:
                return HumanMessage(content=f"[CSV : {file_name} — erreur : {e}]\n\n{text}")

        # ── Fallback (texte brut, code, etc.) ────────────────────────────
        else:
            preview = file_content[:8000] + ("..." if len(file_content) > 8000 else "")
            return HumanMessage(content=f"[Fichier : {file_name}]\n```\n{preview}\n```\n\n{text}")

    def _build_system_prompt(self, user_id: str, deep_research: bool = False, minimal: bool = False, include_kb: bool = True) -> str:
        """Construit le prompt système avec contexte utilisateur.
        - minimal=True : prompt léger sans DB ni KB (mode chat rapide initial)
        - include_kb=False : garde mémoire + préférences mais skip l'appel Firebase KB
          (utilisé pour les workers qui ont leurs propres outils de recherche KB)
        """
        from datetime import date as _date

        if minimal:
            system = SUPERVISOR_RUNTIME_PROMPT.replace(
                "{current_date}", _date.today().strftime("%d/%m/%Y")
            ).replace("{core_memory}", "[Mémoire disponible via outils]").replace("{user_persona}", "")
            return system

        core_memory.set_user_id(user_id)
        core_mem_str = core_memory.format_for_prompt()

        user_persona = ""
        try:
            from core.database import SessionLocal, SystemSetting
            db = SessionLocal()
            pref = db.query(SystemSetting).filter(
                SystemSetting.key == "user_preferences", SystemSetting.user_id == user_id
            ).first()
            if pref and pref.value and pref.value.strip():
                user_persona = f"\n## Préférences Personnelles\n{pref.value.strip()}"
            db.close()
        except Exception as e:
            logger.warning(f"Erreur préférences: {e}")

        knowledge_ctx = ""
        # L'appel Firebase KB est coûteux en latence (~1-3s réseau).
        # Les workers ont leurs propres outils search_knowledge_base → skip.
        if include_kb:
            try:
                from core.firebase_client import knowledge_manager as kb
                docs = kb.get_user_knowledge(user_id)
                if docs:
                    titles = [f"- {d.get('metadata', {}).get('title', '?')}" for d in docs[:10]]
                    knowledge_ctx = "\n\n## Base de Connaissances\n" + "\n".join(titles) + "\nUtilise search_knowledge_base pour le contenu."
            except Exception as e:
                logger.warning(f"Erreur KB: {e}")

        system = SUPERVISOR_RUNTIME_PROMPT.replace(
            "{current_date}", _date.today().strftime("%d/%m/%Y")
        ).replace("{core_memory}", core_mem_str).replace("{user_persona}", user_persona) + knowledge_ctx

        if deep_research:
            system += "\nMODE DEEP RESEARCH: Recherche EXHAUSTIVE obligatoire."

        return system

    async def chat_stream(
        self,
        user_message: str,
        session_id: str,
        file_content: Optional[str] = None,
        file_type: Optional[str] = None,
        file_name: Optional[str] = None,
        user_id: str = "anonymous",
        deep_research: bool = False,
        model_override: Optional[str] = None,
        quality: Optional[str] = "HD",
    ):
        """Streaming : route vers le worker adapté ou répond directement."""
        # Initialiser le contexte
        current_user_id.set(user_id)
        current_session_id.set(session_id)
        req_id = str(uuid.uuid4())
        current_request_id.set(req_id)
        GLOBAL_IMAGE_REGISTRY[req_id] = []
        generated_images.set([])
        knowledge_saved.set([])
        core_memory.set_user_id(user_id)
        current_inline_file.set({
            "content": file_content,
            "type": file_type,
            "name": file_name,
        } if file_content and file_name else None)

        import time
        start_time = time.time()
        first_token_time = None # Added
        has_sent_final = False # Added
        memory = SessionMemory(session_id, user_id=user_id)

        # ── Quota (basé sur le coût réel quotidien) ──────────────────────
        try:
            from core.database import check_quota
            quota = check_quota(user_id, kind="message")
            if not quota["allowed"]:
                plan_name = quota.get("plan_name", "Free")
                reason = quota.get("reason", "budget")
                if reason == "budget":
                    budget = quota.get("daily_budget", 0)
                    msg = f"🔒 Limite quotidienne du plan {plan_name} atteinte (${budget}/jour). Renouvellement à minuit UTC."
                else:
                    msg = f"🔒 Limite du plan {plan_name} atteinte pour aujourd'hui."
                yield {"type": "final", "response": msg, "images": [], "sources": [], "knowledge_saved": [], "suggestions": ["Passer au plan Pro"], "session_id": session_id, "response_time": 0}
                return
        except Exception as e:
            logger.warning(f"[subscription] Erreur: {e}")

        effective_model = model_override

        # ── Classification (0 token, <1ms) ──────────────────────────────────
        worker_type = classify_request(user_message, model_override=effective_model,
                                       file_type=file_type, file_name=file_name)

        # Forcer le bon worker pour les refinements document / image
        try:
            _ref_data = json.loads(user_message)
            if isinstance(_ref_data, dict):
                _selection = _ref_data.get('selection', {})
                _selection_type = _selection.get('type')
                if _selection_type == 'document_refinement':
                    _file_url = _selection.get('file_url', '')
                    _file_name = _selection.get('file_name', '')
                    _target = _file_url or _file_name
                    if _target.endswith('.xlsx') or 'spreadsheet' in _target or 'sheet' in _target.lower():
                        worker_type = 'sheet'
                    else:
                        worker_type = 'document'
                    logger.info(f"[Supervisor] Refinement override -> worker='{worker_type}' pour '{_file_name}'")
                elif _selection_type == 'image_refinement':
                    worker_type = 'media'
                    logger.info("[Supervisor] Image refinement override -> worker='media'")
        except (json.JSONDecodeError, TypeError, KeyError):
            pass

        # Le mode "deep_research" force le chercheur, SAUF si l'utilisateur veut explicitement générer un média.
        if deep_research and worker_type != "media":
            worker_type = "researcher"

        logger.info(f"[Supervisor] Route: '{worker_type}' pour message='{user_message[:60]}...'")

        # ── Config pour les workers ────────────────────────────────────────
        config = {
            "configurable": {
                "user_id": user_id,
                "session_id": session_id,
                "request_id": req_id,
                "model_override": effective_model,
                "quality": quality,
            }
        }

        import asyncio

        human_msg = self._build_human_message(user_message, file_content, file_type, file_name)
        
        # Enregistrer l'attachement utilisateur s'il existe
        user_attachment = None
        if file_content and file_name:
            user_attachment = {
                "kind": "image" if file_type and "image" in file_type else "document",
                "name": file_name,
                "type": file_type,
                "url": None, # L'URL sera gérée par le frontend ou worker si besoin
            }
        
        # Sauvegarde synchrone (très rapide)
        memory.save_message("user", user_message, attachment=user_attachment)

        if not file_content and _is_pure_greeting(user_message):
            final_response, suggestions = _build_greeting_response()
            memory.save_message("assistant", final_response, response_time=0.0)
            yield {
                "type": "final",
                "response": final_response,
                "images": [],
                "sources": [],
                "knowledge_saved": [],
                "suggestions": suggestions,
                "session_id": session_id,
                "response_time": 0.0,
            }
            return

        final_response = ""
        collected_sources = []
        tools_called = []

        # Cache user_email pour record_usage (évite 1 query DB par event worker)
        _cached_user_email = None
        try:
            from core.database import SessionLocal as _SL, UserSubscription as _US
            _db_cache = _SL()
            _sub = _db_cache.query(_US).filter(_US.user_id == user_id).first()
            _cached_user_email = _sub.user_email if _sub and hasattr(_sub, 'user_email') else None
            _db_cache.close()
        except Exception:
            pass

        # ── Route vers le worker ou réponse directe avec TTFB Immédiat ─────
        if worker_type == "chat":
            # Conversation directe — le Supervisor répond lui-même avec VRAI STREAMING
            yield {"type": "thought", "content": "💬 Connexion établie..."}

            # Chargements BDD en parallèle (minimal=True pour éviter Firebase KB ~1-3s)
            history, system_prompt = await asyncio.gather(
                asyncio.to_thread(memory.load_messages, CHAT_SHORT_TERM_MEMORY_MESSAGES),
                asyncio.to_thread(self._build_system_prompt, user_id, deep_research, True, False),
            )

        else:
            worker_map = {
                "researcher": ("researcher", self.researcher),
                "media": ("media", self.media),
                "workspace": ("workspace", self.workspace),
                "document": ("document", self.document),
                "sheet": ("sheet", self.sheet),
            }
            worker_name, worker = worker_map[worker_type]

            yield {"type": "thought", "content": f"Traitement par le module {worker_name}..."}

            if worker_type == "media":
                # Média (image/vidéo) : pas d'historique, pas de system_prompt lourd
                # → zéro attente DB, lancement immédiat
                history = []
                system_prompt = ""
            else:
                # Autres workers : minimal=True pour tous (skip Firebase KB ~1-3s)
                # Les workers ont leurs propres outils search_knowledge_base
                history, system_prompt = await asyncio.gather(
                    asyncio.to_thread(memory.load_messages, CHAT_SHORT_TERM_MEMORY_MESSAGES),
                    asyncio.to_thread(self._build_system_prompt, user_id, deep_research, True, False),
                )

        # Sanitize history
        sanitized = []
        for msg in history:
            if isinstance(msg, AIMessage) and (not msg.content or msg.content == ""):
                msg = AIMessage(content="(action effectuée)")
            elif isinstance(msg, HumanMessage) and (not msg.content or msg.content == ""):
                msg = HumanMessage(content="(message)")
            sanitized.append(msg)
        history = sanitized[-20:]  # Max 20 messages

        if worker_type == "chat":
            if history:
                sys_content = system_prompt
                if len(history) > 10:
                    sys_content = system_prompt[:3000] + "\n[Contexte déjà établi]"
                messages = [
                    HumanMessage(content=f"[Instructions]\n{sys_content}"),
                    AIMessage(content="OK."),
                ] + history + [human_msg]
            else:
                combined = f"[Instructions]\n{system_prompt}\n[Fin]\n\n{human_msg.content}"
                if isinstance(human_msg.content, list):
                    combined = [{"type": "text", "text": f"[Instructions]\n{system_prompt}\n[Fin]\n\n"}] + human_msg.content
                messages = [HumanMessage(content=combined)]

            try:
                # Utilisation de astream pour un retour immédiat des tokens
                llm = get_llm(temperature=0.7, model_override=effective_model, streaming=True) if effective_model else get_llm(temperature=0.7, streaming=True)
                
                full_content = ""
                # Buffer pour intercepter les [SUGGESTION: ...] pendant le streaming
                _stream_buffer = ""

                async for chunk in llm.astream(messages):
                    content = getattr(chunk, "content", "")
                    if content:
                        if first_token_time is None:
                            first_token_time = time.time()
                        if isinstance(content, list):
                            content = "".join([c.get("text", "") if isinstance(c, dict) else str(c) for c in content])
                        full_content += content

                        # Accumulate into buffer and flush only safe content
                        _stream_buffer += content

                        # Check if buffer contains a complete or partial [SUGGESTION: ...] tag
                        while _stream_buffer:
                            bracket_pos = _stream_buffer.find("[")
                            if bracket_pos == -1:
                                # No bracket — flush everything
                                yield {"type": "delta", "content": _stream_buffer}
                                _stream_buffer = ""
                            elif bracket_pos > 0:
                                # Flush text before the bracket
                                yield {"type": "delta", "content": _stream_buffer[:bracket_pos]}
                                _stream_buffer = _stream_buffer[bracket_pos:]
                            else:
                                # Buffer starts with '[' — check if it's a SUGGESTION tag
                                close_pos = _stream_buffer.find("]")
                                if close_pos != -1:
                                    tag = _stream_buffer[:close_pos + 1]
                                    if re.match(r'\[SUGGESTION:', tag, re.IGNORECASE):
                                        # Silently consume the tag
                                        _stream_buffer = _stream_buffer[close_pos + 1:]
                                    else:
                                        # Not a suggestion tag — flush it
                                        yield {"type": "delta", "content": tag}
                                        _stream_buffer = _stream_buffer[close_pos + 1:]
                                else:
                                    # Incomplete bracket — might be a partial SUGGESTION tag, wait for more
                                    if len(_stream_buffer) > 60:
                                        # Too long to be a tag, flush it
                                        yield {"type": "delta", "content": _stream_buffer}
                                        _stream_buffer = ""
                                    break  # Wait for more tokens

                # Flush any remaining buffer (strip suggestions robustement)
                if _stream_buffer:
                    _flushed = re.sub(r'\[SUGGESTION:\s*[^\]]*\]', '', _stream_buffer, flags=re.IGNORECASE).strip()
                    if _flushed:
                        yield {"type": "delta", "content": _flushed}

                final_response = full_content
            except Exception as e:
                logger.error(f"[Supervisor] Erreur chat direct: {e}")
                final_response = "Une erreur est survenue pendant la génération. Veuillez réessayer."

        else:
            # Inclure le contexte de conversation dans la tâche du worker
            history_context = ""
            if history:
                recent = history[-CHAT_SHORT_TERM_MEMORY_MESSAGES:]  # Derniers messages du contexte court
                history_parts = []
                for msg in recent:
                    role = "Utilisateur" if isinstance(msg, HumanMessage) else "FLARE AI"
                    content = msg.content if isinstance(msg.content, str) else str(msg.content)
                    history_parts.append(f"{role}: {content[:300]}")
                history_context = "\n\nContexte de conversation récent:\n" + "\n".join(history_parts)

            # Pour le worker média : injecter le MEDIA_SYSTEM_PROMPT (pas chargé via _build_system_prompt)
            if worker_type == "media":
                from agents.workers.media import MEDIA_SYSTEM_PROMPT as _MSP
                task = f"[Instructions]\n{_MSP}\n[Fin instructions]\n\nRequête utilisateur: {user_message}"
            else:
                task = f"{system_prompt}\n{history_context}\n\nRequête utilisateur: {user_message}"

            # ── CRITIQUE: Injecter le contenu EXTRAIT du fichier dans le message worker ──
            # human_msg contient déjà le contenu extrait (PDF, DOCX, XLSX, image, etc.)
            # On l'ajoute au contexte du worker pour qu'il puisse analyser le fichier
            if worker_type == "media":
                latest_image = memory.get_latest_media_reference(["image"])
                lower_user_message = user_message.lower()
                needs_recent_image = any(token in lower_user_message for token in (
                    "anime", "animer", "animation", "mets-la en mouvement", "met-la en mouvement",
                    "mettre en mouvement", "faire bouger", "donne vie", "change le fond",
                    "fond transparent", "supprime le fond", "enleve le fond", "retire le fond",
                    "transforme cette image", "transforme la photo",
                ))
                if latest_image and needs_recent_image:
                    seeded_media = {
                        "prompt": latest_image.get("name") or "recent image",
                        "type": latest_image.get("type") or "image/jpeg",
                        "name": latest_image.get("name") or "",
                        "url": latest_image.get("url"),
                        "ephemeral": not bool(latest_image.get("url")),
                        "seed_only": True,
                    }
                    if latest_image.get("data"):
                        seeded_media["data"] = latest_image["data"]

                    existing_media = GLOBAL_IMAGE_REGISTRY.get(req_id, []) or []
                    has_seed = any(
                        isinstance(item, dict)
                        and item.get("type", "").startswith("image/")
                        and (item.get("url") == seeded_media.get("url") or item.get("data") == seeded_media.get("data"))
                        for item in existing_media
                    )
                    if not has_seed:
                        existing_media.append(seeded_media)
                        GLOBAL_IMAGE_REGISTRY[req_id] = existing_media
                        generated_images.set(existing_media)

                    task += (
                        "\n\n[CONTEXTE MEDIA RECENT]\n"
                        f"Derniere image disponible: {latest_image.get('name') or 'image'}\n"
                        f"Type: {latest_image.get('type') or 'image/jpeg'}\n"
                        "Une image recente est deja disponible dans le contexte de conversation. "
                        "Si l'utilisateur dit 'anime-la', 'anime cette image', 'mets-la en mouvement' ou fait reference a l'image precedente, "
                        "utilise directement cette image recente sans redemander de lien."
                    )
                    if latest_image.get("url"):
                        task += f"\nURL image recente: {latest_image['url']}"

            if file_content and file_name:
                file_msg = self._build_human_message("", file_content, file_type, file_name)
                file_context = file_msg.content if isinstance(file_msg.content, str) else ""
                # Pour les messages multimodaux (images/vidéos), extraire le texte descriptif
                if isinstance(file_msg.content, list):
                    file_context = f"[Fichier multimodal joint: {file_name} (type: {file_type})]"
                if file_context:
                    task += f"\n\n{file_context}"

            yield {"type": "thought", "content": f"Traitement par le module {worker_name}..."}

            try:
                # Exécuter le worker avec streaming des événements
                worker_graph = worker.graph

                # Préparer les messages pour le worker
                # Pour les images/vidéos en mode chat avec fichier, utiliser human_msg multimodal
                if file_content and file_type and (file_type.startswith("image/") or file_type.startswith("video/")):
                    # Le worker reçoit le message multimodal (image/vidéo inline)
                    worker_messages = [HumanMessage(content=task), human_msg]
                else:
                    worker_messages = [HumanMessage(content=task)]

                async for event in worker_graph.astream_events(
                    {"messages": worker_messages}, config=config, version="v2"
                ):
                    kind = event["event"]
                    if kind == "on_tool_start":
                        tool_name = event["name"]
                        tool_input = event.get("data", {}).get("input", {})
                        tools_called.append(tool_name)
                        display_name = TOOL_DISPLAY_NAMES.get(tool_name, f"Utilisation de {tool_name}...")

                        if tool_name == "web_search" and isinstance(tool_input, dict) and tool_input.get("query"):
                            display_name = f"Recherche : {tool_input['query'][:80]}"
                        elif tool_name == "execute_deep_research" and isinstance(tool_input, dict) and tool_input.get("query"):
                            display_name = f"Recherche approfondie : {tool_input['query'][:60]}"
                        elif tool_name == "search_knowledge_base" and isinstance(tool_input, dict) and tool_input.get("query"):
                            display_name = f"Recherche dans vos documents : {tool_input['query'][:60]}"

                        yield {"type": "thought", "content": display_name}

                    elif kind == "on_tool_end":
                        tool_name = event.get("name", "")
                        if tool_name in ("web_search", "execute_deep_research"):
                            tool_output = event.get("data", {}).get("output", "")
                            output_str = tool_output.content if hasattr(tool_output, "content") else str(tool_output or "")
                            from urllib.parse import urlparse
                            # Extraire les sources
                            for match in re.finditer(r'\[([^\]]+)\]\((https?://[^\)]+)\)', output_str):
                                title, url = match.group(1), match.group(2)
                                url = re.sub(r'[\n\r]+.*$', '', url).rstrip('.,;:!?)>]}')
                                # Ignorer les fichiers internes GCS — ce ne sont pas des sources web
                                if "storage.googleapis.com" in url:
                                    continue
                                try:
                                    domain = urlparse(url).netloc.replace("www.", "")
                                    if domain and len(domain) > 3 and domain not in [s.get("domain") for s in collected_sources]:
                                        collected_sources.append({"url": url, "domain": domain, "title": title[:40]})
                                except Exception:
                                    pass
                            if collected_sources:
                                yield {"type": "thought", "content": f"{len(collected_sources)} source(s) trouvée(s)"}

                    elif kind == "on_chat_model_end":
                        output = event["data"].get("output")
                        if output:
                            has_tool_calls = hasattr(output, "tool_calls") and output.tool_calls
                            if not has_tool_calls:
                                final_response = getattr(output, "content", "")

            except Exception as e:
                logger.error(f"[Supervisor] Erreur worker {worker_name}: {e}", exc_info=True)
                final_response = f"Une erreur est survenue lors de l'exécution du worker {worker_name}: {e}"

        # ── Nettoyage et finalisation ──────────────────────────────────────
        final_response = _clean_response(final_response)

        if not final_response:
            if "generate_image" in tools_called:
                final_response = "Voici l'image que j'ai générée pour vous !"
            elif "generate_video" in tools_called:
                final_response = "Voici la vidéo que j'ai générée pour vous !"
            elif tools_called:
                final_response = f"J'ai utilisé {len(tools_called)} outil(s) pour traiter votre demande."
            else:
                final_response = "Je n'ai pas pu traiter cette demande. Veuillez reformuler."

        # Extraire les sources de la réponse finale (uniquement sources web réelles, pas les fichiers GCS)
        from urllib.parse import urlparse as _urlparse
        for match in re.finditer(r'\[([^\]]+)\]\((https?://[^\)]+)\)', final_response):
            title, url = match.group(1), match.group(2)
            url = re.sub(r'[\n\r]+.*$', '', url).rstrip('.,;:!?)>]}')
            # Ignorer les fichiers GCS internes — ce sont des documents générés, pas des sources
            if "storage.googleapis.com" in url:
                continue
            try:
                domain = _urlparse(url).netloc.replace("www.", "")
                if domain and len(domain) > 3 and domain not in [s.get("domain") for s in collected_sources]:
                    collected_sources.append({"url": url, "domain": domain, "title": title[:40]})
            except Exception:
                pass

        # Images/vidéos générées
        res_images = GLOBAL_IMAGE_REGISTRY.get(req_id, [])
        if not res_images:
            res_images = generated_images.get() or []

        if worker_type == "media" and not res_images and _looks_like_invalid_media_markdown(final_response):
            final_response = (
                "Je n'ai pas pu générer un média exploitable pour cette demande. "
                "Je peux animer l'image actuelle, changer son fond, ou générer une nouvelle vue à partir d'une description précise."
            )

        stream_images = []
        assistant_attachment = None

        for idx, img in enumerate(res_images):
            if img.get("seed_only"):
                continue
            media_obj = {
                "prompt": img.get("prompt", ""), "type": img.get("type", "image/jpeg"),
                "url": img.get("url"), "name": img.get("name"), "ephemeral": img.get("ephemeral", False),
            }
            # Toujours inclure la base64 si disponible : permet l'affichage immédiat sans CORS
            if img.get("data"):
                media_obj["data"] = img["data"]
            
            stream_images.append(media_obj)
            
            # Utiliser le premier média comme attachment principal pour la persistence SQL
            if idx == 0 and not media_obj.get("ephemeral"):
                is_video = media_obj["type"].startswith("video/")
                is_doc = "wordprocessingml" in media_obj["type"] or "document" in media_obj["type"]
                is_sheet = "spreadsheetml" in media_obj["type"] or "excel" in media_obj["type"]
                
                kind = "sheet" if is_sheet else ("document" if is_doc else ("video" if is_video else "image"))
                default_name = "gen.xlsx" if is_sheet else ("gen.docx" if is_doc else ("gen.mp4" if is_video else "gen.jpg"))
                
                assistant_attachment = {
                    "kind": kind,
                    "name": media_obj["name"] or default_name,
                    "type": media_obj["type"],
                    "url": media_obj["url"]
                }
                if kind in {"document", "sheet"} and media_obj.get("data") and not media_obj.get("url"):
                    assistant_attachment["data"] = media_obj["data"]

        # Utiliser d'abord le first token time s'il existe (Fast Path direct Chat),
        # sinon l'utiliser si on a au moins une réponse,
        # sinon prendre le temps actuel en fallback (cas d'erreur)
        ttfb = (first_token_time - start_time) if first_token_time else (time.time() - start_time)
        
        memory.save_message("assistant", final_response, attachment=assistant_attachment, response_time=ttfb)

        # ── Enregistrement usage dans le ledger admin ─────────────────────
        # Les workers loguent déjà les coûts de génération → ne pas doubler
        _has_gen_tools = any(t in tools_called for t in (
            "generate_image", "generate_video", "animate_image",
            "generate_word_document", "generate_excel_document",
        ))
        if not _has_gen_tools:
            try:
                from core.database import record_usage
                # Modèle réel selon le worker
                if worker_type == "chat":
                    _model_for_usage = effective_model or settings.GEMINI_PRO_MODEL
                else:
                    _model_for_usage = effective_model or "gemini-2.5-flash"
                # Différencier deep_research
                _action = "deep_research" if deep_research else worker_type
                record_usage(
                    user_id=user_id,
                    model_name=_model_for_usage,
                    action_kind=_action,
                    user_email=_cached_user_email,
                )
            except Exception as _ue:
                logger.warning(f"[Supervisor] Erreur record_usage: {_ue}")

        # Extraction mémoire automatique — tous les 4 messages
        # La mémoire ne doit pas attendre un 4e message pour devenir utile.
        # On lance l'extraction/synthèse après chaque échange assistant.
        asyncio.create_task(self._run_background_memory_tasks(memory, user_id))

        # Extraire les suggestions de la réponse finale
        final_response, suggestions = _extract_suggestions(final_response)
        suggestions = _ensure_suggestions(suggestions, worker_type)

        yield {
            "type": "final",
            "response": final_response,
            "images": stream_images,
            "sources": collected_sources[:15],
            "knowledge_saved": knowledge_saved.get() or [],
            "suggestions": suggestions,
            "session_id": session_id,
            "response_time": float(ttfb)
        }

        asyncio.create_task(self._cleanup_registry(req_id))

    async def chat(
        self,
        user_message: str,
        session_id: str,
        file_content: Optional[str] = None,
        file_type: Optional[str] = None,
        file_name: Optional[str] = None,
        user_id: str = "anonymous",
        deep_research: bool = False,
        model_override: Optional[str] = None,
        quality: Optional[str] = "HD",
    ) -> dict:
        """Version non-streaming (fallback)."""
        final = None
        async for event in self.chat_stream(
            user_message=user_message, session_id=session_id,
            file_content=file_content, file_type=file_type, file_name=file_name,
            user_id=user_id, deep_research=deep_research, model_override=model_override,
            quality=quality,
        ):
            if isinstance(event, dict) and event.get("type") == "final":
                final = event

        if final:
            return {
                "response": final["response"],
                "images": final.get("images", []),
                "knowledge_saved": final.get("knowledge_saved", []),
                "suggestions": final.get("suggestions", []),
            }
        return {"response": "Erreur de traitement.", "images": [], "knowledge_saved": [], "suggestions": []}

    async def _run_background_memory_tasks(self, memory: SessionMemory, user_id: str):
        try:
            from core.memory import CoreMemory
            local_core = CoreMemory(user_id)
            all_messages = memory.load_messages()
            bg_llm = get_llm(temperature=0.3, model_override=settings.GEMINI_ROUTING_MODEL)
            await local_core.auto_extract_facts(all_messages, bg_llm)
            await memory.summarize_if_needed(bg_llm)
            logger.info(f"[background_memory] Extraction mémoire terminée pour {user_id}")
        except Exception as e:
            logger.warning(f"[background_memory] Erreur: {e}")

    async def _cleanup_registry(self, req_id: str):
        await asyncio.sleep(60)
        GLOBAL_IMAGE_REGISTRY.pop(req_id, None)


# ─── Instance singleton ──────────────────────────────────────────────────────

_supervisor_instance: SupervisorAgent | None = None


def get_supervisor() -> SupervisorAgent:
    """Retourne l'instance singleton du Supervisor."""
    global _supervisor_instance
    if _supervisor_instance is None:
        _supervisor_instance = SupervisorAgent()
    return _supervisor_instance







